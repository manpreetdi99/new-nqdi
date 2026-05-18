from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.util as util

# ─── COLOUR PALETTE ───────────────────────────────────────────────────────────
DARK_BG    = RGBColor(0x0F, 0x17, 0x2A)   # very dark navy
CARD_BG    = RGBColor(0x1A, 0x25, 0x3D)   # dark blue-grey (card surface)
ACCENT     = RGBColor(0x00, 0xB4, 0xD8)   # cyan accent
ACCENT2    = RGBColor(0x48, 0xCA, 0xE4)   # lighter cyan
GREEN      = RGBColor(0x06, 0xD6, 0xA0)   # teal/green
ORANGE     = RGBColor(0xFF, 0xB7, 0x03)   # amber
RED        = RGBColor(0xEF, 0x47, 0x6F)   # red
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xB0, 0xBE, 0xD4)
MID_GREY   = RGBColor(0x64, 0x74, 0x8B)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_layout(prs):
    return prs.slide_layouts[6]          # completely blank


def fill_bg(slide, colour=DARK_BG):
    """Fill slide background with a solid colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = colour


def add_rect(slide, l, t, w, h, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.fill.background()          # no border
    return shape


def add_text(slide, text, l, t, w, h,
             size=18, bold=False, colour=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = colour
    return txBox


def add_bullet_box(slide, items, l, t, w, h,
                   title=None, title_size=14, item_size=12,
                   title_colour=ACCENT, item_colour=LIGHT_GREY,
                   bg_colour=CARD_BG, padding=0.18):
    """Draw a rounded card with optional title + bullet list."""
    add_rect(slide, l, t, w, h, bg_colour)
    y = t + padding
    if title:
        add_text(slide, title, l + padding, y, w - 2*padding, 0.38,
                 size=title_size, bold=True, colour=title_colour)
        y += 0.42
    for item in items:
        add_text(slide, f"  {item}", l + padding, y, w - 2*padding, 0.33,
                 size=item_size, colour=item_colour)
        y += 0.3


# ══════════════════════════════════════════════════════════════════════════════
#   SLIDE BUILDERS
# ══════════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    """Slide 1 – Title / Cover"""
    sl = prs.slides.add_slide(blank_layout(prs))
    fill_bg(sl)

    # gradient-like top bar
    add_rect(sl, 0, 0, 13.33, 0.12, ACCENT)

    # big glowing circle (decorative)
    circ = sl.shapes.add_shape(9, Inches(9.5), Inches(1.2), Inches(5.5), Inches(5.5))
    circ.fill.solid()
    circ.fill.fore_color.rgb = RGBColor(0x00, 0x4E, 0x64)
    circ.line.fill.background()

    # product name chip
    add_rect(sl, 0.6, 1.4, 2.8, 0.42, ACCENT)
    add_text(sl, "FASMETRICS ANALYTICS", 0.6, 1.4, 2.8, 0.42,
             size=11, bold=True, colour=DARK_BG, align=PP_ALIGN.CENTER)

    # main title
    add_text(sl, "Performance Insights", 0.6, 1.95, 8.5, 1.2,
             size=54, bold=True, colour=WHITE)

    # subtitle
    add_text(sl, "Telecom Network Quality & SQL Benchmarking Platform",
             0.6, 3.25, 8.5, 0.7, size=22, colour=ACCENT2)

    # description
    add_text(sl, "A full-stack analytics dashboard for telecom engineers\n"
                 "to benchmark SQL Server queries and analyse\n"
                 "mobile network KPIs in real time.",
             0.6, 4.05, 7.5, 1.2, size=16, colour=LIGHT_GREY, italic=True)

    # bottom bar with tech tags
    add_rect(sl, 0, 6.9, 13.33, 0.6, CARD_BG)
    tags = ["React 18", "TypeScript", "FastAPI", "SQL Server", "Recharts", "Leaflet", "Tailwind CSS"]
    x = 0.5
    for tag in tags:
        w = len(tag) * 0.115 + 0.25
        add_rect(sl, x, 6.97, w, 0.36, MID_GREY)
        add_text(sl, tag, x + 0.04, 6.97, w - 0.08, 0.36,
                 size=11, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)
        x += w + 0.18

    # bottom right date
    add_text(sl, "May 2026", 11.5, 7.05, 1.6, 0.35,
             size=13, colour=MID_GREY, align=PP_ALIGN.RIGHT)

    return sl


def slide_overview(prs):
    """Slide 2 – Project Overview"""
    sl = prs.slides.add_slide(blank_layout(prs))
    fill_bg(sl)
    add_rect(sl, 0, 0, 13.33, 0.08, ACCENT)

    add_text(sl, "Project Overview", 0.5, 0.2, 12, 0.65,
             size=32, bold=True, colour=WHITE)
    add_rect(sl, 0.5, 0.88, 2.2, 0.05, ACCENT)

    # LEFT column – What / Who / Why
    cards = [
        ("WHAT",   ACCENT,  ["Full-stack web application",
                              "React + TypeScript frontend",
                              "Python FastAPI backend",
                              "SQL Server database integration"]),
        ("WHO",    GREEN,   ["Telecom network engineers",
                              "COSMOTE / Greek operator teams",
                              "Performance & QA analysts"]),
        ("WHY",    ORANGE,  ["Benchmark slow SQL queries",
                              "Analyse call quality (MOS, setup time)",
                              "Visualise 4G/5G antenna coverage",
                              "Identify dropped & failed calls"]),
    ]
    y_start = 1.15
    for label, colour, items in cards:
        add_rect(sl, 0.5, y_start, 5.5, len(items)*0.32 + 0.65, CARD_BG)
        add_text(sl, label, 0.68, y_start + 0.1, 1.5, 0.38,
                 size=13, bold=True, colour=colour)
        for i, item in enumerate(items):
            add_text(sl, f"• {item}", 0.68, y_start + 0.5 + i*0.31,
                     5.1, 0.3, size=12.5, colour=LIGHT_GREY)
        y_start += len(items)*0.32 + 0.75

    # RIGHT column – architecture diagram (text-based)
    add_rect(sl, 6.8, 1.15, 6.0, 5.8, CARD_BG)
    add_text(sl, "Architecture", 7.0, 1.28, 5.6, 0.45,
             size=16, bold=True, colour=ACCENT)

    boxes = [
        (7.8, 1.85, 4.0, 0.55, RGBColor(0x17, 0x3F, 0x5F), "BROWSER", "React 18 + TypeScript", ACCENT2),
        (7.8, 2.65, 4.0, 0.55, RGBColor(0x0D, 0x3B, 0x2E), "BACKEND", "FastAPI (Python 3.11)", GREEN),
        (7.8, 3.45, 4.0, 0.55, RGBColor(0x3B, 0x1F, 0x00), "DATABASE", "SQL Server + ODBC", ORANGE),
        (7.8, 4.25, 4.0, 0.55, RGBColor(0x28, 0x10, 0x10), "FILES", "Excel (geo4g.xlsx)", RED),
        (7.8, 5.05, 4.0, 0.55, RGBColor(0x1A, 0x1A, 0x3A), "MAPS", "Leaflet + OpenStreetMap", ACCENT),
    ]
    for bx, by, bw, bh, bc, lbl, sub, lc in boxes:
        add_rect(sl, bx, by, bw, bh, bc)
        add_text(sl, lbl, bx + 0.15, by + 0.02, 1.3, 0.28,
                 size=10, bold=True, colour=lc)
        add_text(sl, sub, bx + 0.15, by + 0.27, 3.5, 0.25,
                 size=11.5, colour=WHITE)
        # connector arrow
        if by < 5.05:
            add_text(sl, "↓", bx + 1.85, by + 0.56, 0.4, 0.3,
                     size=14, colour=MID_GREY, align=PP_ALIGN.CENTER)

    return sl


def slide_features(prs):
    """Slide 3 – Key Features"""
    sl = prs.slides.add_slide(blank_layout(prs))
    fill_bg(sl)
    add_rect(sl, 0, 0, 13.33, 0.08, ACCENT)

    add_text(sl, "Key Features", 0.5, 0.2, 12, 0.65,
             size=32, bold=True, colour=WHITE)
    add_rect(sl, 0.5, 0.88, 1.9, 0.05, ACCENT)

    features = [
        ("SQL Benchmarking",     ACCENT,   Inches(0.5),  Inches(1.1),
         ["• Batch execution of SQL queries",
          "• Execution time & row count metrics",
          "• Stats: total time, avg time, query count",
          "• Pre-built query template library",
          "• Smart chart auto-selection"]),
        ("All Calls Analysis",   GREEN,    Inches(3.6),  Inches(1.1),
         ["• Filter by collection, location, status",
          "• Color-coded call statuses",
          "• Setup time, duration, MOS display",
          "• Click to open detailed call view",
          "• Comment/annotation support"]),
        ("Data Sessions",        ORANGE,   Inches(6.7),  Inches(1.1),
         ["• Test type & direction filters",
          "• Throughput & ping metrics",
          "• YouTube MOS scoring",
          "• Side-by-side A/B device view",
          "• Per-session timeline events"]),
        ("Antennas Map",         RED,      Inches(9.8),  Inches(1.1),
         ["• 4G/5G cell tower database",
          "• Leaflet interactive map",
          "• Site ID, Cell ID, PCI, Azimuth",
          "• Downtilt, height, vendor info",
          "• Greece coverage visualization"]),
    ]

    for title, colour, lx, ly, bullets in features:
        card = sl.shapes.add_shape(1, lx, ly, Inches(2.8), Inches(5.5))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.fill.background()

        top_bar = sl.shapes.add_shape(1, lx, ly, Inches(2.8), Inches(0.08))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = colour
        top_bar.line.fill.background()

        txb = sl.shapes.add_textbox(lx + Inches(0.15), ly + Inches(0.15),
                                    Inches(2.5), Inches(0.45))
        tf = txb.text_frame
        p  = tf.paragraphs[0]
        r  = p.add_run()
        r.text = title
        r.font.size  = Pt(15)
        r.font.bold  = True
        r.font.color.rgb = colour

        for i, bullet in enumerate(bullets):
            txb2 = sl.shapes.add_textbox(lx + Inches(0.15),
                                         ly + Inches(0.7) + Inches(i * 0.42),
                                         Inches(2.55), Inches(0.4))
            tf2 = txb2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            r2 = p2.add_run()
            r2.text = bullet
            r2.font.size  = Pt(12)
            r2.font.color.rgb = LIGHT_GREY

    return sl


def slide_sql_benchmark(prs):
    """Slide 4 – SQL Benchmarking Deep-Dive"""
    sl = prs.slides.add_slide(blank_layout(prs))
    fill_bg(sl)
    add_rect(sl, 0, 0, 13.33, 0.08, ACCENT)

    add_text(sl, "SQL Benchmarking Engine", 0.5, 0.2, 12, 0.65,
             size=32, bold=True, colour=WHITE)
    add_rect(sl, 0.5, 0.88, 3.0, 0.05, ACCENT)

    # workflow steps
    steps = [
        ("1", "Select Database", "Choose SQL Server DB from dropdown",     ACCENT),
        ("2", "Write / Build SQL", "SQL mode or visual Query Builder UI",  ACCENT2),
        ("3", "Run Benchmark", "POST /api/benchmark — batch execution",    GREEN),
        ("4", "View Stats", "Time, rows, avg time per query",              ORANGE),
        ("5", "Analyse Charts", "Auto-selected Bar / Area / Radar chart",  RED),
    ]

    for i, (num, title, desc, col) in enumerate(steps):
        x = 0.5 + i * 2.55
        add_rect(sl, x, 1.1, 2.35, 1.7, CARD_BG)
        # circle number
        circ = sl.shapes.add_shape(9, Inches(x + 0.85), Inches(1.2),
                                   Inches(0.65), Inches(0.65))
        circ.fill.solid()
        circ.fill.fore_color.rgb = col
        circ.line.fill.background()
        add_text(sl, num, x + 0.85, 1.2, 0.65, 0.65,
                 size=18, bold=True, colour=DARK_BG, align=PP_ALIGN.CENTER)
        add_text(sl, title, x + 0.12, 1.98, 2.1, 0.38,
                 size=13, bold=True, colour=col, align=PP_ALIGN.CENTER)
        add_text(sl, desc, x + 0.12, 2.4, 2.1, 0.35,
                 size=11, colour=LIGHT_GREY, align=PP_ALIGN.CENTER)

        if i < 4:
            add_text(sl, "→", x + 2.35, 1.78, 0.25, 0.4,
                     size=18, colour=MID_GREY, align=PP_ALIGN.CENTER)

    # API payload box
    add_rect(sl, 0.5, 3.05, 5.8, 4.0, CARD_BG)
    add_text(sl, "API Request  (POST /api/benchmark)", 0.65, 3.15, 5.5, 0.4,
             size=14, bold=True, colour=ACCENT)
    code_lines = [
        '{',
        '  "database": "COSMOTE_LTE_2025",',
        '  "queries": [',
        '    "SELECT TOP 200 * FROM CallAnalysis",',
        '    "SELECT AVG(MOS) FROM Calls WHERE tech=\'LTE\'"',
        '  ]',
        '}',
    ]
    for i, line in enumerate(code_lines):
        add_text(sl, line, 0.65, 3.6 + i * 0.42, 5.3, 0.4,
                 size=12, colour=GREEN, italic=True)

    # Response metrics box
    add_rect(sl, 6.7, 3.05, 6.15, 4.0, CARD_BG)
    add_text(sl, "Response Metrics", 6.85, 3.15, 5.8, 0.4,
             size=14, bold=True, colour=ORANGE)
    metrics = [
        ("Query Label",      "Descriptive name per query",   ACCENT2),
        ("Execution Time",   "Milliseconds per query",       GREEN),
        ("Rows Returned",    "COUNT(*) from result set",     ORANGE),
        ("Column Names",     "Schema metadata",              ACCENT),
        ("Data Array",       "Full result rows as JSON",     LIGHT_GREY),
    ]
    for i, (k, v, c) in enumerate(metrics):
        y = 3.65 + i * 0.6
        add_rect(sl, 6.85, y, 1.8, 0.42, RGBColor(0x11, 0x1F, 0x38))
        add_text(sl, k, 6.9, y + 0.04, 1.7, 0.35, size=11, bold=True, colour=c)
        add_text(sl, v, 8.75, y + 0.04, 3.8, 0.35, size=11, colour=LIGHT_GREY)

    return sl


def slide_telecom(prs):
    """Slide 5 – Telecom Analytics"""
    sl = prs.slides.add_slide(blank_layout(prs))
    fill_bg(sl)
    add_rect(sl, 0, 0, 13.33, 0.08, GREEN)

    add_text(sl, "Telecom Network Analytics", 0.5, 0.2, 12, 0.65,
             size=32, bold=True, colour=WHITE)
    add_rect(sl, 0.5, 0.88, 3.5, 0.05, GREEN)

    # KPI boxes top row
    kpis = [
        ("MOS Score",    "Mean Opinion Score\nVoice quality 1–5",            GREEN),
        ("Setup Time",   "Call setup latency\nin milliseconds",              ACCENT),
        ("RSRP / RSRQ",  "LTE signal strength\nand quality indicators",     ORANGE),
        ("Throughput",   "DL/UL speeds\nMbps per session",                  RED),
        ("Call Status",  "Completed / Dropped\nFailed / System Release",    ACCENT2),
    ]
    for i, (k, v, c) in enumerate(kpis):
        x = 0.4 + i * 2.55
        add_rect(sl, x, 1.1, 2.38, 1.3, CARD_BG)
        add_rect(sl, x, 1.1, 2.38, 0.07, c)
        add_text(sl, k, x + 0.12, 1.22, 2.1, 0.38, size=13, bold=True, colour=c)
        add_text(sl, v, x + 0.12, 1.63, 2.1, 0.7, size=11.5, colour=LIGHT_GREY)

    # Two columns below
    # Left – Call Record fields
    add_bullet_box(sl,
        ["Session ID (unique call identifier)",
         "Technology: 5G SA / 5G NSA / 4G LTE / 3G",
         "Call Type: Voice, Video, VoLTE, VoNR",
         "Duration (seconds), Setup Time (ms)",
         "Avg MOS (1–5 scale)",
         "Location, Latitude / Longitude",
         "isValid flag + comment annotation",
         "A-side & B-side device measurement"],
        l=0.4, t=2.6, w=5.9, h=4.5,
        title="Call Record Fields",
        title_size=14, item_size=12.5,
        title_colour=GREEN)

    # Right – Technology breakdown
    add_rect(sl, 6.8, 2.6, 6.15, 4.5, CARD_BG)
    add_text(sl, "Network Technologies Tracked", 6.95, 2.7, 5.8, 0.45,
             size=14, bold=True, colour=ACCENT)

    techs = [
        ("5G SA",      "Standalone 5G NR",              RGBColor(0x06, 0xD6, 0xA0)),
        ("5G NSA",     "5G New Radio + 4G anchor",      RGBColor(0x00, 0xB4, 0xD8)),
        ("4G LTE",     "Long Term Evolution",           RGBColor(0xFF, 0xB7, 0x03)),
        ("3G UMTS",    "Universal Mobile Telecom",      RGBColor(0xEF, 0x47, 0x6F)),
        ("2G GSM",     "Global System for Mobile",      RGBColor(0xA8, 0xDA, 0xDC)),
        ("VoLTE",      "Voice over LTE (HD Voice)",     RGBColor(0xFF, 0xD1, 0x66)),
        ("VoNR",       "Voice over New Radio (5G)",     RGBColor(0x90, 0xE0, 0xEF)),
    ]
    for i, (tech, desc, col) in enumerate(techs):
        y = 3.25 + i * 0.51
        add_rect(sl, 6.95, y, 1.2, 0.38, col)
        add_text(sl, tech, 6.95, y, 1.2, 0.38,
                 size=12, bold=True, colour=DARK_BG, align=PP_ALIGN.CENTER)
        add_text(sl, desc, 8.25, y + 0.03, 4.5, 0.32, size=12, colour=LIGHT_GREY)

    return sl


def slide_architecture(prs):
    """Slide 6 – Tech Stack & Architecture"""
    sl = prs.slides.add_slide(blank_layout(prs))
    fill_bg(sl)
    add_rect(sl, 0, 0, 13.33, 0.08, ORANGE)

    add_text(sl, "Technology Stack", 0.5, 0.2, 12, 0.65,
             size=32, bold=True, colour=WHITE)
    add_rect(sl, 0.5, 0.88, 2.5, 0.05, ORANGE)

    # Frontend stack
    add_rect(sl, 0.4, 1.1, 5.9, 6.0, CARD_BG)
    add_rect(sl, 0.4, 1.1, 5.9, 0.07, ACCENT)
    add_text(sl, "Frontend", 0.6, 1.18, 5.5, 0.45,
             size=16, bold=True, colour=ACCENT)

    fe_items = [
        ("React 18",          "UI framework + hooks",          ACCENT),
        ("TypeScript",        "Type-safe development",         ACCENT2),
        ("Vite 7",            "Build tool & dev server",       GREEN),
        ("Tailwind CSS 3.4",  "Utility-first styling",         ORANGE),
        ("shadcn/ui",         "Radix UI component library",    LIGHT_GREY),
        ("Recharts",          "Bar, Area, Radar charts",       RED),
        ("Leaflet",           "Interactive map tiles",         GREEN),
        ("React Router v6",   "Client-side routing",           ACCENT2),
        ("Framer Motion",     "Page & component animations",   ORANGE),
        ("TanStack Query",    "Server state management",       ACCENT),
        ("React Hook Form",   "Form management + Zod",         LIGHT_GREY),
        ("Lucide React",      "SVG icon library",              MID_GREY),
    ]
    for i, (name, desc, col) in enumerate(fe_items):
        y = 1.72 + i * 0.38
        add_text(sl, f"• {name}", 0.6, y, 2.3, 0.35, size=12, bold=True, colour=col)
        add_text(sl, desc, 3.0, y, 3.1, 0.35, size=11.5, colour=LIGHT_GREY)

    # Backend stack
    add_rect(sl, 6.9, 1.1, 6.05, 6.0, CARD_BG)
    add_rect(sl, 6.9, 1.1, 6.05, 0.07, GREEN)
    add_text(sl, "Backend", 7.1, 1.18, 5.7, 0.45,
             size=16, bold=True, colour=GREEN)

    be_items = [
        ("FastAPI",          "Async Python REST framework",   GREEN),
        ("Python 3.11",      "Core runtime",                  LIGHT_GREY),
        ("pyodbc",           "SQL Server ODBC driver",        ORANGE),
        ("ODBC Driver 17",   "Microsoft SQL connectivity",    ORANGE),
        ("SQLAlchemy",       "ORM (minimal usage)",           LIGHT_GREY),
        ("openpyxl",         "Excel file reader (antennas)",  ACCENT),
        ("python-dotenv",    "Environment config / secrets",  MID_GREY),
        ("Uvicorn",          "ASGI server",                   GREEN),
        ("CORS Middleware",  "Cross-origin request support",  LIGHT_GREY),
        ("SQL Server",       "Production database",           RED),
        ("Excel .xlsx",      "Antenna geo-database",          ORANGE),
        ("LAN Deployment",   "192.168.x.x multi-user access", ACCENT2),
    ]
    for i, (name, desc, col) in enumerate(be_items):
        y = 1.72 + i * 0.38
        add_text(sl, f"• {name}", 7.1, y, 2.4, 0.35, size=12, bold=True, colour=col)
        add_text(sl, desc, 9.55, y, 3.2, 0.35, size=11.5, colour=LIGHT_GREY)

    return sl


def slide_api(prs):
    """Slide 7 – API Endpoints"""
    sl = prs.slides.add_slide(blank_layout(prs))
    fill_bg(sl)
    add_rect(sl, 0, 0, 13.33, 0.08, RED)

    add_text(sl, "REST API Endpoints", 0.5, 0.2, 12, 0.65,
             size=32, bold=True, colour=WHITE)
    add_rect(sl, 0.5, 0.88, 2.8, 0.05, RED)

    categories = [
        ("Databases & Collections", ACCENT, [
            "GET  /api/databases",
            "GET  /api/collections?database=X",
            "GET  /api/locations?database=X&collection=Y",
        ]),
        ("Call & Session Data", GREEN, [
            "GET  /api/calls  (with filters)",
            "GET  /api/data_calls",
            "POST /api/calls/comment",
        ]),
        ("Performance", ORANGE, [
            "POST /api/benchmark",
            "GET  /api/results_kpi",
        ]),
        ("LTE Measurements", ACCENT2, [
            "GET  /api/lte_values",
            "GET  /api/lte_values_b_side",
            "GET  /api/lte_measurement_comparison",
            "GET  /api/lte_scanner_measurement",
            "GET  /api/lte_scanner_raw",
        ]),
        ("GSM Measurements", RGBColor(0xA8, 0xDA, 0xDC), [
            "GET  /api/gsm_values",
            "GET  /api/gsm_values_b_side",
            "GET  /api/gsm_context_signal",
            "GET  /api/gsm_scanner_raw",
        ]),
        ("Call Context & Analysis", RED, [
            "GET  /api/call_context_signal",
            "GET  /api/call_context_technology",
            "GET  /api/call_side_comparison",
            "GET  /api/call_paging_info",
            "GET  /api/call_device_info",
            "GET  /api/tracelog_values",
            "GET  /api/cell_info",
        ]),
        ("Infrastructure", ORANGE, [
            "GET  /api/antennas",
        ]),
    ]

    positions = [
        (0.4,  1.1,  3.95, 4.5),
        (0.4,  5.75, 3.95, 1.8),
        (4.7,  1.1,  4.0,  2.1),
        (4.7,  3.3,  4.0,  2.85),
        (4.7,  6.25, 4.0,  2.4),
        (9.0,  1.1,  4.0,  4.0),
        (9.0,  5.25, 4.0,  1.1),
    ]

    for (title, col, eps), (lx, ly, lw, lh) in zip(categories, positions):
        add_rect(sl, lx, ly, lw, lh, CARD_BG)
        add_rect(sl, lx, ly, lw, 0.06, col)
        add_text(sl, title, lx+0.12, ly+0.1, lw-0.2, 0.38,
                 size=12, bold=True, colour=col)
        for i, ep in enumerate(eps):
            add_text(sl, ep, lx+0.12, ly+0.52+i*0.37, lw-0.2, 0.34,
                     size=10.5, colour=LIGHT_GREY, italic=True)

    # total count badge
    add_rect(sl, 4.95, 0.22, 1.55, 0.45, RED)
    add_text(sl, "25+ Endpoints", 4.95, 0.22, 1.55, 0.45,
             size=12, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)

    return sl


def slide_ui_ux(prs):
    """Slide 8 – UI/UX Features"""
    sl = prs.slides.add_slide(blank_layout(prs))
    fill_bg(sl)
    add_rect(sl, 0, 0, 13.33, 0.08, ACCENT2)

    add_text(sl, "UI / UX Highlights", 0.5, 0.2, 12, 0.65,
             size=32, bold=True, colour=WHITE)
    add_rect(sl, 0.5, 0.88, 2.6, 0.05, ACCENT2)

    ui_features = [
        ("Dark Theme",         "Navy/cyan dark palette — built for long monitoring sessions", ACCENT),
        ("Responsive Design",  "Tailwind CSS — works on desktop, tablet, and mobile screens", GREEN),
        ("Toast Notifications","Sonner library — success/error/info toasts on every action", ORANGE),
        ("Smart Charts",       "Auto-detects data shape → picks Bar, Area, or Radar chart",  RED),
        ("Persistent State",   "localStorage remembers DB, collection, tab, location filters", ACCENT2),
        ("Scroll-to-View",     "Auto-scrolls to selected call row when switching between tabs", GREEN),
        ("Framer Motion",      "Smooth page transitions and card entrance animations",          ORANGE),
        ("Color-coded Status", "Red=Invalid, Orange=Dropped, Violet=System Release, Green=OK", RED),
        ("Dual Editor Modes",  "Raw SQL editor OR visual Query Builder — same output",         ACCENT),
        ("Greek Locale",       "Number formatting (1.234,56) for Greek telecom standards",    LIGHT_GREY),
    ]

    cols = 2
    per_col = (len(ui_features) + 1) // 2
    for i, (title, desc, col) in enumerate(ui_features):
        col_idx = i // per_col
        row_idx = i % per_col
        x = 0.4 + col_idx * 6.45
        y = 1.1 + row_idx * 0.6
        add_rect(sl, x, y, 6.1, 0.52, CARD_BG)
        add_rect(sl, x, y, 0.06, 0.52, col)
        add_text(sl, title, x + 0.18, y + 0.04, 2.0, 0.42,
                 size=13, bold=True, colour=col)
        add_text(sl, desc, x + 2.25, y + 0.08, 3.7, 0.36, size=11.5, colour=LIGHT_GREY)

    # screenshot placeholder
    add_rect(sl, 0.4, 7.05, 12.55, 0.35, CARD_BG)
    add_text(sl, "Application available on LAN: http://192.168.10.44:5173  |  API: http://192.168.10.44:8000",
             0.55, 7.07, 12.3, 0.3, size=11, colour=MID_GREY, align=PP_ALIGN.CENTER)

    return sl


def slide_data_flow(prs):
    """Slide 9 – Data Flow Diagram"""
    sl = prs.slides.add_slide(blank_layout(prs))
    fill_bg(sl)
    add_rect(sl, 0, 0, 13.33, 0.08, ACCENT)

    add_text(sl, "Data Flow", 0.5, 0.2, 12, 0.65,
             size=32, bold=True, colour=WHITE)
    add_rect(sl, 0.5, 0.88, 1.6, 0.05, ACCENT)

    # Data flow nodes
    nodes = [
        (0.5,  2.4, 2.4, 1.1, CARD_BG,   ACCENT, "ENGINEER",        "Selects DB, enters SQL,\napplies filters"),
        (3.5,  2.4, 2.4, 1.1, CARD_BG,   GREEN,  "REACT APP",       "Manages state, calls\nFastAPI endpoints"),
        (6.5,  2.4, 2.4, 1.1, CARD_BG,   ORANGE, "FASTAPI",         "Validates params,\nexecutes queries"),
        (9.5,  2.4, 2.4, 1.1, CARD_BG,   RED,    "SQL SERVER",      "Returns rows,\ntimings, metadata"),
        (6.5,  4.7, 2.4, 1.1, CARD_BG,   ACCENT2,"EXCEL FILE",      "geo4g.xlsx\nAntenna locations"),
        (3.5,  4.7, 2.4, 1.1, CARD_BG,   GREEN,  "LEAFLET MAP",     "Renders cell towers\non OSM tiles"),
        (0.5,  4.7, 2.4, 1.1, CARD_BG,   ACCENT, "RECHARTS",        "Renders Bar/Area/\nRadar charts"),
    ]

    for nx, ny, nw, nh, nbg, nc, label, desc in nodes:
        add_rect(sl, nx, ny, nw, nh, nbg)
        add_rect(sl, nx, ny, nw, 0.06, nc)
        add_text(sl, label, nx+0.1, ny+0.12, nw-0.2, 0.38, size=12, bold=True, colour=nc)
        add_text(sl, desc,  nx+0.1, ny+0.52, nw-0.2, 0.52, size=10.5, colour=LIGHT_GREY)

    # Arrows top row
    arrows_h = [(2.95, 2.88, 0.5, 0.32, "→"), (5.95, 2.88, 0.5, 0.32, "→"), (8.95, 2.88, 0.5, 0.32, "→")]
    for ax, ay, aw, ah, sym in arrows_h:
        add_text(sl, sym, ax, ay, aw, ah, size=20, colour=MID_GREY, align=PP_ALIGN.CENTER)

    # Vertical arrows
    add_text(sl, "↓", 7.6, 3.56, 0.4, 0.4, size=20, colour=MID_GREY, align=PP_ALIGN.CENTER)
    add_text(sl, "←", 5.95, 5.17, 0.5, 0.32, size=20, colour=MID_GREY, align=PP_ALIGN.CENTER)
    add_text(sl, "←", 2.95, 5.17, 0.5, 0.32, size=20, colour=MID_GREY, align=PP_ALIGN.CENTER)

    # Return arrow label
    add_text(sl, "JSON results\n(rows, times)", 5.95, 3.6, 0.5, 0.9, size=9, colour=MID_GREY)

    # Key insight box
    add_rect(sl, 0.4, 6.2, 12.55, 1.1, CARD_BG)
    add_text(sl, "Deployment Note:", 0.6, 6.3, 2.5, 0.38, size=13, bold=True, colour=ACCENT)
    add_text(sl, "Frontend and backend both run on the same LAN server. "
                 "Multiple engineers can connect simultaneously from their own machines. "
                 "API base URL is configured via VITE_API_BASE_URL environment variable. "
                 "SQL Server credentials are stored in backend/.env file.",
             0.6, 6.68, 12.1, 0.55, size=12, colour=LIGHT_GREY)

    return sl


def slide_project_structure(prs):
    """Slide 10 – Project Structure"""
    sl = prs.slides.add_slide(blank_layout(prs))
    fill_bg(sl)
    add_rect(sl, 0, 0, 13.33, 0.08, GREEN)

    add_text(sl, "Project Structure", 0.5, 0.2, 12, 0.65,
             size=32, bold=True, colour=WHITE)
    add_rect(sl, 0.5, 0.88, 2.6, 0.05, GREEN)

    # Frontend files
    add_rect(sl, 0.4, 1.1, 5.9, 5.95, CARD_BG)
    add_rect(sl, 0.4, 1.1, 5.9, 0.07, ACCENT)
    add_text(sl, "Frontend  (src/)", 0.58, 1.18, 5.6, 0.42, size=14, bold=True, colour=ACCENT)

    fe_files = [
        ("pages/Index.tsx",              "Main dashboard, 1000+ lines",         ACCENT),
        ("components/QueryEditor.tsx",   "SQL editor with dual modes",           ACCENT2),
        ("components/BenchmarkCharts.tsx","Recharts visualization",              GREEN),
        ("components/CallsList.tsx",     "Call records table + filters",         ORANGE),
        ("components/CallDetail.tsx",    "Detailed call view",                   RED),
        ("components/AntennasMap.tsx",   "Leaflet cell tower map",               GREEN),
        ("components/DataSessionsList.tsx","Data session table",                 ACCENT2),
        ("lib/api.ts",                   "HTTP client, 500+ lines",              ACCENT),
        ("lib/callData.ts",              "Mock data generation",                 MID_GREY),
        ("types/benchmark.ts",           "TypeScript interfaces",                ORANGE),
        ("hooks/use-local-storage.ts",   "Persistent state hook",                LIGHT_GREY),
        ("components/ui/",              "30+ shadcn/ui components",             MID_GREY),
    ]
    for i, (name, desc, col) in enumerate(fe_files):
        y = 1.65 + i * 0.38
        add_text(sl, name, 0.58, y, 2.8, 0.35, size=10.5, bold=True, colour=col)
        add_text(sl, desc, 3.45, y, 2.7, 0.35, size=10.5, colour=LIGHT_GREY)

    # Backend files
    add_rect(sl, 6.8, 1.1, 6.15, 3.2, CARD_BG)
    add_rect(sl, 6.8, 1.1, 6.15, 0.07, GREEN)
    add_text(sl, "Backend  (backend/)", 6.98, 1.18, 5.8, 0.42, size=14, bold=True, colour=GREEN)

    be_files = [
        ("app.py",          "FastAPI routes, 2154 lines",     GREEN),
        ("db.py",           "SQL Server connection helpers",  ACCENT),
        ("queries.py",      "SQL templates & helpers",        ORANGE),
        ("requirements.txt","Python dependencies",            MID_GREY),
        (".env",            "DB credentials (not committed)", RED),
    ]
    for i, (name, desc, col) in enumerate(be_files):
        y = 1.65 + i * 0.42
        add_text(sl, name, 6.98, y, 2.3, 0.35, size=11.5, bold=True, colour=col)
        add_text(sl, desc, 9.35, y, 3.4, 0.35, size=11.5, colour=LIGHT_GREY)

    # Config files
    add_rect(sl, 6.8, 4.45, 6.15, 2.6, CARD_BG)
    add_rect(sl, 6.8, 4.45, 6.15, 0.07, ORANGE)
    add_text(sl, "Config & Data", 6.98, 4.53, 5.8, 0.42, size=14, bold=True, colour=ORANGE)

    cfg_files = [
        ("vite.config.ts",       "Vite & proxy settings",             ACCENT),
        ("tailwind.config.ts",   "Tailwind theme customization",      GREEN),
        ("package.json",         "~60 npm dependencies",              ORANGE),
        ("tsconfig.app.json",    "TypeScript compiler config",        LIGHT_GREY),
        ("geo4g.xlsx",           "4G/5G antenna database",            RED),
    ]
    for i, (name, desc, col) in enumerate(cfg_files):
        y = 4.98 + i * 0.38
        add_text(sl, name, 6.98, y, 2.6, 0.35, size=11, bold=True, colour=col)
        add_text(sl, desc, 9.65, y, 3.1, 0.35, size=11, colour=LIGHT_GREY)

    return sl


def slide_summary(prs):
    """Slide 11 – Summary & Conclusion"""
    sl = prs.slides.add_slide(blank_layout(prs))
    fill_bg(sl)
    add_rect(sl, 0, 0, 13.33, 0.08, ACCENT)

    add_text(sl, "Summary", 0.5, 0.2, 12, 0.65, size=32, bold=True, colour=WHITE)
    add_rect(sl, 0.5, 0.88, 1.5, 0.05, ACCENT)

    # 3-column summary
    cols_data = [
        ("What it does",   ACCENT,  [
            "SQL benchmarking with timing",
            "Voice call quality analytics",
            "Data session KPI analysis",
            "4G/5G antenna mapping",
            "Multi-collection filtering",
            "Real-time LTE/GSM metrics",
        ]),
        ("Built with",     GREEN,   [
            "React 18 + TypeScript",
            "FastAPI + Python 3.11",
            "SQL Server via pyodbc",
            "Recharts data visualization",
            "Leaflet interactive maps",
            "Tailwind + shadcn/ui",
        ]),
        ("For whom",       ORANGE,  [
            "COSMOTE network engineers",
            "QA & performance analysts",
            "Drive-test teams",
            "Database administrators",
            "Network planning teams",
            "Operations center staff",
        ]),
    ]
    for i, (title, col, items) in enumerate(cols_data):
        x = 0.4 + i * 4.25
        add_rect(sl, x, 1.1, 4.0, 5.5, CARD_BG)
        add_rect(sl, x, 1.1, 4.0, 0.07, col)
        add_text(sl, title, x+0.15, 1.2, 3.7, 0.45, size=16, bold=True, colour=col)
        for j, item in enumerate(items):
            add_text(sl, f"✓  {item}", x+0.15, 1.75 + j*0.65, 3.7, 0.55,
                     size=13, colour=LIGHT_GREY)

    # Bottom banner
    add_rect(sl, 0.4, 6.85, 12.55, 0.52, CARD_BG)
    add_text(sl, "Performance Insights  |  FASMETRICS Analytics  |  Telecom Network Quality Platform  |  May 2026",
             0.5, 6.88, 12.3, 0.42, size=13, colour=MID_GREY, align=PP_ALIGN.CENTER)

    return sl


# ══════════════════════════════════════════════════════════════════════════════
#   MAIN
# ══════════════════════════════════════════════════════════════════════════════
if __name__ == "__main__":
    prs = new_prs()

    slide_title(prs)
    slide_overview(prs)
    slide_features(prs)
    slide_sql_benchmark(prs)
    slide_telecom(prs)
    slide_architecture(prs)
    slide_api(prs)
    slide_ui_ux(prs)
    slide_data_flow(prs)
    slide_project_structure(prs)
    slide_summary(prs)

    out = r"c:\Users\Mechanical Engineer\Documents\performance-insights\FASMETRICS_Analytics_Presentation.pptx"
    prs.save(out)
    print(f"Saved: {out}")
    print(f"Slides: {len(prs.slides)}")
